#!/usr/bin/env python3
"""지구환경과 생태 강의 PPTX → 텍스트 추출."""
from __future__ import annotations

from pathlib import Path

from pptx import Presentation

ROOT = Path(__file__).resolve().parents[1]
OUT = ROOT / ".tmp-pdf-text"
BASE = Path(r"c:\Users\tmddb\Desktop\과제\26-1학기\지구과학\강의자료 모음")

PAIRS = [
    (BASE / "3주차_지구환경과생태학_근대.pptx", "eco-week03-modern.txt"),
    (BASE / "4주차강의록_지구환경과생태.pptx", "eco-week04-env-eco.txt"),
    (BASE / "5주차_피터싱어종차별주의.pptx", "eco-week05-singer.txt"),
    (BASE / "6주차_T레건_동물권리론.pptx", "eco-week06-regan.txt"),
    (BASE / "7주차_종차별주의 옹호논변들.pptx", "eco-week07-speciesism.txt"),
    (BASE / "9주차_슈바이처_생명외경 (1).pptx", "eco-week09-schweitzer.txt"),
    (BASE / "12주차_테일러의생명중심주의.pptx", "eco-week12-taylor.txt"),
    (BASE / "13주차_레오폴드_대지윤리.pptx", "eco-week13-leopold.txt"),
    (BASE / "14주차_아르네 네스_심층생태학.pptx", "eco-week14-naess.txt"),
]

OUT.mkdir(parents=True, exist_ok=True)


def slide_text(slide) -> str:
    parts: list[str] = []
    for shape in slide.shapes:
        if not hasattr(shape, "text"):
            continue
        t = (shape.text or "").strip()
        if t:
            parts.append(t)
    return "\n".join(parts)


for src, name in PAIRS:
    if not src.exists():
        print(f"MISSING: {src}")
        continue
    prs = Presentation(str(src))
    parts = []
    for i, slide in enumerate(prs.slides):
        text = slide_text(slide)
        parts.append(f"--- slide {i + 1} ---\n{text}")
    dst = OUT / name
    dst.write_text("\n\n".join(parts), encoding="utf-8")
    print(name, len(prs.slides), dst.stat().st_size)
