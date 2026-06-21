#!/usr/bin/env python3
"""인공지능 기초 기말 PPTX → 텍스트 추출."""
from __future__ import annotations

from pathlib import Path

from pptx import Presentation

ROOT = Path(__file__).resolve().parents[1]
OUT = ROOT / ".tmp-pdf-text"
BASE = Path(r"c:\Users\tmddb\Desktop\과제\26-1학기\인공지능기초\기말 강의자료")

PAIRS = [
    (BASE / "07_인공지능과 데이터 과학.pptx", "ai-week07-data-science.txt"),
    (BASE / "08_스스로 학습하는 머신러닝 (1).pptx", "ai-week08-ml.txt"),
    (BASE / "09_인공 신경망 기술 (1).pptx", "ai-week09-ann.txt"),
    (BASE / "10_딥러닝의 세계로.pptx", "ai-week10-dl.txt"),
    (BASE / "11_대규모 언어 모델과 생성형 인공지능.pptx", "ai-week11-llm.txt"),
    (BASE / "12_인간과 겨루는 인공지능- 강화 학습_S.pptx", "ai-week12-rl.txt"),
]

OUT.mkdir(parents=True, exist_ok=True)


def slide_text(slide) -> str:
    parts: list[str] = []
    for shape in slide.shapes:
        if not hasattr(shape, "text"):
            continue
        t = (shape.text or "").strip()
        if t:
            parts.append(t)
    return "\n".join(parts)


for src, name in PAIRS:
    if not src.exists():
        print(f"MISSING: {src}")
        continue
    prs = Presentation(str(src))
    parts = []
    for i, slide in enumerate(prs.slides):
        text = slide_text(slide)
        parts.append(f"--- slide {i + 1} ---\n{text}")
    dst = OUT / name
    dst.write_text("\n\n".join(parts), encoding="utf-8")
    print(name, len(prs.slides), dst.stat().st_size)
